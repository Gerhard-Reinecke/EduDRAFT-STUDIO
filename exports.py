# ======================================================================================
# exports.py
# ======================================================================================
# Module: Export Rendering, Visual Packaging & Output Assembly Engine
#
# System: EduDraft Studio (Marike App)
# Version: 4.4.10
#
# --------------------------------------------------------------------------------------
# PURPOSE
# --------------------------------------------------------------------------------------
# This module is the final output assembly layer of EduDraft Studio.
#
# It converts structured document content into export-ready formats, manages visual
# placeholder rendering, packages diagrams/charts/images into DOCX and PPTX outputs,
# handles LaTeX rendering for presentation workflows, and builds preview-safe HTML for
# in-app document review.
#
# It acts as the bridge between structured draft content and user-deliverable files.
#
# --------------------------------------------------------------------------------------
# CORE RESPONSIBILITIES
# --------------------------------------------------------------------------------------
# 1. DOCX Export Assembly
#    - Converts Markdown-based worksheet content into DOCX using Pandoc
#    - Resolves and embeds generated visual assets into the export package
#    - Strips non-DOCX bundle sections such as PPT outlines and answer-key wrappers
#    - Preserves editable equations where supported
#
# 2. PPTX Export Assembly
#    - Converts structured slide outlines into PowerPoint presentations
#    - Places bullet content, equations, diagrams, charts, and generated images onto slides
#    - Keeps visuals constrained within slide content boxes
#    - Supports speaker notes integration
#
# 3. Visual Placeholder Rendering
#    - Parses [[VISUAL ...]] placeholders and related visual markers
#    - Routes:
#        • charts → chart renderer
#        • diagrams / graphs → diagram_library
#        • images / pictures → image generation renderer
#        • unknown visual types → placeholder visual generator
#
# 4. Chart & Diagram Packaging
#    - Generates PNG outputs for charts from inline datasets
#    - Produces placeholder visuals when data is missing or malformed
#    - Integrates rendered diagrams into both document and presentation exports
#
# 5. LaTeX Rendering for Presentation Output
#    - Renders LaTeX expressions into PNGs for PPT slide placement
#    - Supports inline and display-style equation handling
#    - Falls back safely with error-image generation when rendering fails
#
# 6. Preview HTML Construction
#    - Builds MathJax-compatible HTML previews for in-app document review
#    - Escapes raw content safely for browser display
#    - Provides non-export preview output for user validation
#
# --------------------------------------------------------------------------------------
# KEY SUBSYSTEMS
# --------------------------------------------------------------------------------------
# - Bundle splitter for DOCX-safe worksheet extraction
# - VISUAL placeholder parser and renderer
# - Chart generation helpers
# - DOCX export pipeline
# - PPTX slide rendering pipeline
# - LaTeX-to-image conversion helpers
# - MathJax HTML preview builder
#
# --------------------------------------------------------------------------------------
# DEPENDENCIES
# --------------------------------------------------------------------------------------
# - Pandoc → DOCX generation
# - python-pptx → PowerPoint generation
# - matplotlib → charts, placeholders, and equation rendering
# - Pillow (PIL) → image handling
# - latex2mathml / cairosvg → alternate math rendering pathway
# - diagram_library.py → diagram/image generation integration
# - Standard library utilities for file, regex, and subprocess handling
#
# --------------------------------------------------------------------------------------
# DESIGN PRINCIPLES
# --------------------------------------------------------------------------------------
# - Final-output reliability over theoretical purity
# - Graceful fallback when visuals or data are incomplete
# - Clear routing between chart, diagram, image, and placeholder logic
# - Export symmetry across DOCX, PPTX, and preview workflows
# - Safe packaging of temporary assets for downstream conversion tools
#
# --------------------------------------------------------------------------------------
# SYSTEM ROLE
# --------------------------------------------------------------------------------------
# This module is the final packaging and rendering layer of the platform.
#
# Upstream modules produce structured content.
# exports.py transforms that structured content into user-facing deliverables.
#
# In practical terms:
#   - llm.py creates the structured output
#   - template_engine.py shapes layout and presentation intent
#   - exports.py turns that into actual files the user can open, download, and use
#
# --------------------------------------------------------------------------------------
# NOTES
# --------------------------------------------------------------------------------------
# - This module is operationally critical for all export workflows
# - Changes here directly affect DOCX quality, PPTX layout fidelity, and preview accuracy
# - Pandoc availability is mandatory for DOCX export
# - Visual placeholder handling here must remain aligned with llm.py output rules
# - Temporary image asset handling is essential for reliable file packaging
#
# ======================================================================================


import os
import re
import io
import tempfile
import subprocess
import shutil
from typing import Optional, List, Tuple, Dict

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PInches, Pt

from latex2mathml.converter import convert as tex_to_mathml
import cairosvg

import diagram_library
from diagram_library import render_diagram_png


# =============================
# BUNDLE SPLITTER (DOCX FIX)
# =============================
_BUNDLE_MARKER_RE = re.compile(
    r"(?is)\n\s*(?:\d+\)\s*)?(PPT_OUTLINE|ANSWER_KEY)\s*:\s*"
)

def _strip_bundle_sections_for_docx(text: str) -> str:
    """
    DOCX should contain ONLY the worksheet markdown, not the PPT_OUTLINE / ANSWER_KEY bundle.
    If markers exist, truncate at the first marker.
    """
    s = (text or "").strip()
    m = _BUNDLE_MARKER_RE.search(s)
    if not m:
        return s
    return s[:m.start()].rstrip()


# =============================
# VISUAL PLACEHOLDERS
# =============================
_VISUAL_RE = re.compile(r'\[\[VISUAL\s+(.+?)\]\]{1,2}["\']?', re.IGNORECASE)
_BARE_VID_RE = re.compile(r"^\s*v(\d+)\s*$", re.IGNORECASE)

_VISUAL_PLACEHOLDER_RE = re.compile(
    r"📊\s*Visual placeholder\s+(id=.*?)(?=$|\n)",
    re.IGNORECASE
)

_MD_IMG_RE = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
_DIAGRAM_ATTACHED_RE = re.compile(r"\bdiagram\s+attached\b", re.IGNORECASE)


def _maybe_autogen_diagram_png(slide_title: str, slide_lines: List[str]) -> Optional[str]:
    blob = "\n".join([slide_title or ""] + (slide_lines or []))
    if not _DIAGRAM_ATTACHED_RE.search(blob):
        return None

    where = (slide_title or "Diagram").strip()
    prompt = blob.strip()
    notes = ""
    return render_diagram_png(where=where, prompt=prompt, notes=notes)


# =============================
# VISUAL RENDERING HELPERS
# =============================

def _parse_visual_kv(payload: str) -> Dict[str, str]:
    out: Dict[str, str] = {}
    for m in re.finditer(r'(\w+)\s*=\s*"([^"]*)"', payload):
        out[m.group(1).strip().lower()] = m.group(2)
    return out


def _wrap_text_lines(text: str, width: int = 60) -> str:
    words = (text or "").split()
    lines: List[str] = []
    cur: List[str] = []
    n = 0
    for w in words:
        if n + len(w) + (1 if cur else 0) > width:
            lines.append(" ".join(cur))
            cur = [w]
            n = len(w)
        else:
            cur.append(w)
            n += len(w) + (1 if cur else 0)
    if cur:
        lines.append(" ".join(cur))
    return "\n".join(lines)


def _make_placeholder_visual_png(kind: str, where: str, prompt: str, notes: str = "") -> str:
    title = f"{(kind or 'visual').upper()} — {where or 'unspecified'}"
    body = _wrap_text_lines(prompt or "(no prompt)", 70)
    foot = _wrap_text_lines(notes or "", 70)

    fig = plt.figure(figsize=(8, 3.2), dpi=140)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")

    ax.text(0.02, 0.92, title, fontsize=14, weight="bold", va="top")
    ax.text(0.02, 0.78, body, fontsize=11, va="top")

    if foot.strip():
        ax.text(0.02, 0.16, "Notes:", fontsize=10, weight="bold", va="top")
        ax.text(0.02, 0.10, foot, fontsize=10, va="top")

    ax.add_patch(plt.Rectangle((0.01, 0.02), 0.98, 0.96, fill=False, linewidth=1))

    tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(tmp_img.name, format="png", bbox_inches="tight", pad_inches=0.2)
    plt.close(fig)
    return tmp_img.name


def _try_make_chart_png(data: str, title: str = "") -> str:
    d = (data or "").strip()
    if not d or d.upper() == "MISSING":
        return _make_placeholder_visual_png("chart", title, "Chart data is missing.", "")

    # (3) Grouped bars format: "2020: A=3,B=4; 2021: A=2,B=5"
    if ";" in d and ":" in d and "=" in d:
        try:
            year_blocks = [b.strip() for b in d.split(";") if b.strip()]
            years: List[str] = []
            series_names: List[str] = []
            series_values: Dict[str, List[float]] = {}

            for block in year_blocks:
                year, rest = block.split(":", 1)
                year = year.strip()
                if not year:
                    continue

                pairs = [p.strip() for p in rest.split(",") if p.strip()]
                vals_for_year: Dict[str, float] = {}
                for p in pairs:
                    if "=" not in p:
                        continue
                    sname, sval = p.split("=", 1)
                    sname = sname.strip()
                    sval_f = float(sval.strip())
                    vals_for_year[sname] = sval_f
                    if sname not in series_names:
                        series_names.append(sname)

                years.append(year)
                for sname in series_names:
                    series_values.setdefault(sname, [])
                    series_values[sname].append(vals_for_year.get(sname, 0.0))

            if not years or not series_names:
                return _make_placeholder_visual_png("chart", title, f"Could not parse chart data: {d}", "")

            fig = plt.figure(figsize=(8.0, 3.8), dpi=140)
            ax = fig.add_axes([0.08, 0.18, 0.90, 0.72])

            x = list(range(len(years)))
            n = len(series_names)
            total_width = 0.8
            bar_w = total_width / max(n, 1)

            for i, sname in enumerate(series_names):
                offset = (i - (n - 1) / 2) * bar_w
                ax.bar(
                    [xi + offset for xi in x],
                    series_values.get(sname, [0.0] * len(years)),
                    width=bar_w,
                    label=sname,
                )

            ax.set_xticks(x)
            ax.set_xticklabels(years)
            ax.set_xlabel("Year")
            ax.set_ylabel("Value")
            if title:
                ax.set_title(title)
            ax.grid(True, axis="y", linestyle="--", alpha=0.3)
            ax.legend()

            tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            fig.savefig(tmp_img.name, format="png", bbox_inches="tight", pad_inches=0.2)
            plt.close(fig)
            return tmp_img.name
        except Exception:
            return _make_placeholder_visual_png("chart", title, f"Could not parse chart data: {d}", "")

    # (2) Series format: "x: Jan,Feb,Mar | y: 10,12,9"
    if "|" in d and ":" in d:
        try:
            parts = d.split("|")
            x_part = parts[0].split(":", 1)[1]
            y_part = parts[1].split(":", 1)[1]
            xs = [x.strip() for x in x_part.split(",") if x.strip()]
            ys = [float(y.strip()) for y in y_part.split(",") if y.strip()]

            fig = plt.figure(figsize=(7.5, 3.5), dpi=140)
            ax = fig.add_axes([0.10, 0.18, 0.86, 0.72])
            ax.plot(xs, ys, marker="o")
            if title:
                ax.set_title(title)
            ax.grid(True, linestyle="--", alpha=0.3)

            tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            fig.savefig(tmp_img.name, format="png", bbox_inches="tight", pad_inches=0.2)
            plt.close(fig)
            return tmp_img.name
        except Exception:
            return _make_placeholder_visual_png("chart", title, f"Could not parse series data: {d}", "")

    # (1) Categorical bar format: "A:3, B:4, C:2"
    try:
        items: List[Tuple[str, float]] = []
        for part in d.split(","):
            part = part.strip()
            if not part or ":" not in part:
                continue
            k, v = part.split(":", 1)
            items.append((k.strip(), float(v.strip())))

        if not items:
            return _make_placeholder_visual_png("chart", title, f"Could not parse chart data: {d}", "")

        labels = [k for k, _ in items]
        values = [v for _, v in items]

        fig = plt.figure(figsize=(7.5, 3.5), dpi=140)
        ax = fig.add_axes([0.10, 0.22, 0.86, 0.68])
        ax.bar(labels, values)
        if title:
            ax.set_title(title)
        ax.grid(True, axis="y", linestyle="--", alpha=0.3)

        tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        fig.savefig(tmp_img.name, format="png", bbox_inches="tight", pad_inches=0.2)
        plt.close(fig)
        return tmp_img.name
    except Exception:
        return _make_placeholder_visual_png("chart", title, f"Could not parse chart data: {d}", "")


def render_visuals_for_export(markdown_text: str) -> Tuple[str, List[str]]:
    """
    Converts VISUAL placeholders into real images for Pandoc export.
    Returns: (new_markdown, generated_image_paths)
    """
    generated: List[str] = []

    def _normalize_quotes(s: str) -> str:
        return (
            (s or "")
            .replace("“", '"').replace("”", '"')
            .replace("’", "'")
            .replace("–", "-").replace("—", "-")
        )

    def _render_from_payload(payload: str) -> str:
        payload2 = _normalize_quotes(payload)
        kv = _parse_visual_kv(payload2)

        vid = kv.get("id", f"v{len(generated)+1}")
        kind = (kv.get("kind", "visual") or "visual").lower()
        where = kv.get("where", "")
        prompt = kv.get("prompt", "")
        notes = kv.get("notes", "")
        data = kv.get("data", "")

        # ROUTING:
        # chart -> chart renderer (data-based)
        # graph/diagram -> diagram_library (prompt-based)
        # image/photo/picture -> image_gen via diagram_library
        if kind in {"chart"}:
            img_path = _try_make_chart_png(data=data, title=(where or "").strip())
        elif kind in {"graph", "diagram", "table"}:
            img_path = diagram_library.render_diagram_png(
                where=where,
                prompt=prompt,
                notes=(notes + ("\nDATA: " + data if data else "")).strip(),
                subtype=kv.get("subtype")
            )
        elif kind in {"image", "photo", "picture"}:
            img_path = diagram_library.render_image_gen_png(where=where, prompt=prompt, notes=notes)
        else:
            img_path = _make_placeholder_visual_png(kind, where, prompt, notes)

        generated.append(img_path)
        return f'\n\n![{vid}]({img_path})\n\n'

    text = _normalize_quotes(markdown_text or "")

    # (A) Convert 📊 Visual placeholder lines
    def repl_placeholder(m):
        payload = (m.group(1) or "").strip()
        return _render_from_payload(payload)

    text = _VISUAL_PLACEHOLDER_RE.sub(repl_placeholder, text)

    # (B) Convert bare visual IDs lines ("v1")
    lines = text.splitlines()
    out_lines: List[str] = []
    for idx, line in enumerate(lines):
        m = _BARE_VID_RE.match(line)
        if not m:
            out_lines.append(line)
            continue

        vid = f"v{m.group(1)}"
        start = max(0, idx - 12)
        end = min(len(lines), idx + 13)
        context = "\n".join(lines[start:end])

        payload = (
            f'id="{vid}" kind="diagram" where="auto" '
            f'prompt="{context.replace(chr(34), chr(39))}" notes=""'
        )
        out_lines.append(_render_from_payload(payload))

    text = "\n".join(out_lines)

    # (C) Convert [[VISUAL ...]]
    def repl_visual(m):
        payload = (m.group(1) or "").strip()
        return _render_from_payload(payload)

    new_md = _VISUAL_RE.sub(repl_visual, text)
    return new_md, generated


def extract_visuals_from_line(line: str) -> Tuple[str, List[str]]:
    """
    If line contains a [[VISUAL ...]] block, return:
      (clean_text_without_visual, [image_paths])
    Used for PPT bullets.
    """
    imgs: List[str] = []

    def repl(m):
        payload = m.group(1).strip()
        payload = (
            payload
            .replace("“", '"').replace("”", '"')
            .replace("’", "'")
            .replace("–", "-").replace("—", "-")
        )

        kv = _parse_visual_kv(payload)
        kind = (kv.get("kind", "visual") or "visual").lower()
        where = kv.get("where", "")
        prompt = kv.get("prompt", "")
        notes = kv.get("notes", "")
        data = kv.get("data", "")

        if kind in {"chart"}:
            img_path = _try_make_chart_png(data=data, title=(where or "").strip())
        elif kind in {"graph", "diagram", "table"}:
            img_path = diagram_library.render_diagram_png(
                where=where,
                prompt=prompt,
                notes=(notes + ("\nDATA: " + data if data else "")).strip(),
                subtype=kv.get("subtype")
            )
        elif kind in {"image", "photo", "picture"}:
            img_path = diagram_library.render_image_gen_png(where=where, prompt=prompt, notes=notes)
        else:
            img_path = _make_placeholder_visual_png(kind, where, prompt, notes)

        imgs.append(img_path)
        return ""

    cleaned = _VISUAL_RE.sub(repl, line or "").strip()
    return cleaned, imgs


# =============================
# EXPORTS
# =============================

def _strip_outer_md_fence(md_text: str) -> str:
    """
    If the entire document is wrapped in a single fenced code block
    (``` or ```markdown), strip the outer fence so Pandoc can render
    headings/images instead of exporting them as literal text.
    """
    s = (md_text or "").strip()
    if not s.startswith("```"):
        return md_text

    lines = s.splitlines()
    if not lines:
        return md_text

    first = lines[0].strip()
    last = lines[-1].strip()

    if first.startswith("```") and last == "```":
        return "\n".join(lines[1:-1]).strip()

    return md_text


def normalize_math_delimiters_for_pandoc(md_text: str) -> str:
    fixed = md_text.replace(r"\[", "$$").replace(r"\]", "$$")
    fixed = fixed.replace(r"\(", "$").replace(r"\)", "$")
    return fixed


def md_to_docx_with_editable_equations(md_text: str, out_path: str, *, pre_rendered: bool = False) -> None:
    """
    DOCX = worksheet markdown ONLY.
    pre_rendered=True means md_text already contains resolved image links (from render_visuals_for_export),
    so we must NOT call render_visuals_for_export again (prevents double-charge / double-render).
    """
    tmpdir = tempfile.mkdtemp()

    # ---- HARD CHECK: Pandoc must exist
    if not shutil.which("pandoc"):
        raise RuntimeError("Pandoc is not installed or not on PATH. DOCX export requires pandoc in Docker build.")

    md_path = os.path.join(tmpdir, "doc.md")

    md_text = _strip_bundle_sections_for_docx(md_text)
    md_text = _strip_outer_md_fence(md_text)
    fixed = normalize_math_delimiters_for_pandoc(md_text)

    generated_imgs: List[str] = []

    if pre_rendered:
        # Extract any markdown image paths already present
        generated_imgs = [m.group(1).strip() for m in _MD_IMG_RE.finditer(fixed)]
    else:
        fixed, generated_imgs = render_visuals_for_export(fixed)

    # Copy images into tmpdir and rewrite paths to local relative filenames
    for i, src_path in enumerate(generated_imgs or []):
        if not src_path or not os.path.exists(src_path):
            continue
        local_name = f"visual_{i+1}.png"
        dst_path = os.path.join(tmpdir, local_name)
        try:
            with open(src_path, "rb") as fsrc:
                data = fsrc.read()
            with open(dst_path, "wb") as fdst:
                fdst.write(data)
            fixed = fixed.replace(f"({src_path})", f"({local_name})")
        except Exception:
            pass

    with open(md_path, "w", encoding="utf-8") as f:
        f.write(fixed)

    subprocess.run(
        ["pandoc", "doc.md", "-o", out_path, "--resource-path=."],
        check=True,
        cwd=tmpdir
    )


# =============================
# LaTeX rendering for PPT
# =============================

LATEX_INLINE_RE = re.compile(r"\\\((.+?)\\\)")
LATEX_DISPLAY_RE = re.compile(r"\\\[(.+?)\\\]", re.S)


def render_latex_png(latex: str, dpi: int = 220) -> io.BytesIO:
    """
    PPT LaTeX rendering:
    Use matplotlib mathtext (reliable).
    """
    latex2 = (latex or "").strip()
    try:
        return render_latex_matplotlib(latex2, dpi)
    except Exception:
        return create_error_image(f"LaTeX Error: {latex2[:30]}...")


def render_latex_via_mathml(latex: str) -> io.BytesIO:
    mathml = tex_to_mathml(latex)

    svg_template = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" width="1000" height="300">
  <style>
    math {{ font-size: 40px; }}
  </style>
  {mathml}
</svg>'''

    buf = io.BytesIO()
    cairosvg.svg2png(bytestring=svg_template.encode('utf-8'), write_to=buf)
    buf.seek(0)

    img = Image.open(buf)

    max_width = 800
    if img.size[0] > max_width:
        width_percent = max_width / float(img.size[0])
        new_height = int(float(img.size[1]) * float(width_percent))
        img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)

    final_buf = io.BytesIO()
    img.save(final_buf, format='PNG')
    final_buf.seek(0)
    return final_buf


def render_latex_matplotlib(latex: str, dpi: int = 220) -> io.BytesIO:
    expr = f"${latex}$"

    fig = plt.figure(figsize=(0.01, 0.01), dpi=dpi)
    fig.patch.set_alpha(0.0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")

    t = ax.text(0, 0, expr, fontsize=22, va="bottom", ha="left")
    fig.canvas.draw()
    bbox = t.get_window_extent(renderer=fig.canvas.get_renderer())
    plt.close(fig)

    pad_px = 10
    width_in = (bbox.width + pad_px * 2) / dpi
    height_in = (bbox.height + pad_px * 2) / dpi

    fig2 = plt.figure(figsize=(max(width_in, 0.6), max(height_in, 0.3)), dpi=dpi)
    fig2.patch.set_alpha(0.0)
    ax2 = fig2.add_axes([0, 0, 1, 1])
    ax2.axis("off")
    ax2.text(pad_px / dpi, pad_px / dpi, expr, fontsize=22, va="bottom", ha="left")

    buf = io.BytesIO()
    fig2.savefig(buf, format="png", transparent=True, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig2)
    buf.seek(0)
    return buf


def create_error_image(message: str) -> io.BytesIO:
    fig = plt.figure(figsize=(4, 0.5), dpi=100)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, message, fontsize=12, ha='center', va='center', color='red', alpha=0.7)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True, bbox_inches='tight', pad_inches=0.1)
    plt.close(fig)
    buf.seek(0)
    return buf


def split_text_and_math(s: str):
    tokens = []
    i = 0
    while i < len(s):
        m_inline = LATEX_INLINE_RE.search(s, i)
        m_disp = LATEX_DISPLAY_RE.search(s, i)

        m = None
        kind = None
        if m_inline and m_disp:
            if m_inline.start() < m_disp.start():
                m, kind = m_inline, "inline"
            else:
                m, kind = m_disp, "display"
        elif m_inline:
            m, kind = m_inline, "inline"
        elif m_disp:
            m, kind = m_disp, "display"

        if not m:
            tokens.append(("text", s[i:]))
            break

        if m.start() > i:
            tokens.append(("text", s[i:m.start()]))

        tokens.append(("math", m.group(1).strip(), kind))
        i = m.end()

    merged = []
    for tok in tokens:
        if tok[0] == "text":
            if tok[1]:
                if merged and merged[-1][0] == "text":
                    merged[-1] = ("text", merged[-1][1] + tok[1])
                else:
                    merged.append(tok)
        else:
            merged.append(tok)
    return merged


def _add_picture_fit(slide, img_path: str, x, y, max_w, max_h, center: bool = True):
    if max_w <= 0 or max_h <= 0:
        return None

    pic = slide.shapes.add_picture(img_path, x, y)
    w, h = pic.width, pic.height

    scale = min(max_w / w, max_h / h, 1.0)
    new_w = int(w * scale)
    new_h = int(h * scale)

    pic.width = new_w
    pic.height = new_h
    if center:
        pic.left = int(x + (max_w - new_w) / 2)
    pic.top = y
    return pic


def outline_to_pptx_with_math(outline: str, out_path: str) -> None:
    prs = Presentation()
    slides = re.split(r"\n(?=Slide\s+\d+:)", (outline or "").strip(), flags=re.M)

    for chunk in slides:
        chunk = (chunk or "").strip()
        if not chunk:
            continue

        title_match = re.match(r"Slide\s+\d+:\s*(.*)", chunk)
        title = title_match.group(1).strip() if title_match else "Slide"

        notes = ""
        m_notes = re.search(r"Speaker notes:\s*(.*)$", chunk, flags=re.S)
        if m_notes:
            notes = m_notes.group(1).strip()
            chunk_wo_notes = re.sub(r"Speaker notes:\s*.*$", "", chunk, flags=re.S).strip()
        else:
            chunk_wo_notes = chunk

        content_lines = [ln.rstrip() for ln in chunk_wo_notes.splitlines()[1:] if ln.strip()]

        bullet_lines: List[str] = []
        standalone_imgs: List[str] = []

        for ln in content_lines:
            line = ln.strip()
            mimg = _MD_IMG_RE.search(line)
            if mimg:
                standalone_imgs.append(mimg.group(1).strip())
                continue
            if line.startswith("- "):
                bullet_lines.append(line[2:].strip())
            else:
                bullet_lines.append(line)

        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Clean LaTeX in titles so it shows as normal text
        title_clean = LATEX_INLINE_RE.sub(lambda m: m.group(1), title)
        title_clean = LATEX_DISPLAY_RE.sub(lambda m: m.group(1), title_clean)
        slide.shapes.title.text = title_clean

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        content_shape = slide.shapes.placeholders[1]
        box_left = content_shape.left
        box_top = content_shape.top
        box_width = content_shape.width
        box_height = content_shape.height
        y_cursor = box_top + PInches(0.40)

        def remaining_h():
            return (box_top + box_height) - y_cursor - PInches(0.15)

        auto_png = _maybe_autogen_diagram_png(title, content_lines)
        if auto_png:
            pic = _add_picture_fit(
                slide,
                auto_png,
                box_left + PInches(0.25),
                y_cursor,
                box_width - PInches(0.5),
                remaining_h(),
            )
            if pic is not None:
                y_cursor += pic.height + PInches(0.15)

        for img_path in standalone_imgs:
            pic = _add_picture_fit(
                slide,
                img_path,
                box_left + PInches(0.25),
                y_cursor,
                box_width - PInches(0.5),
                remaining_h(),
            )
            if pic is not None:
                y_cursor += pic.height + PInches(0.15)

        for i, b in enumerate(bullet_lines):
            cleaned_b, visual_imgs = extract_visuals_from_line(b)

            # If bullet is purely an image placeholder -> insert image(s) instead of text
            if visual_imgs and not cleaned_b:
                for img_path in visual_imgs:
                    pic = _add_picture_fit(
                        slide,
                        img_path,
                        box_left + PInches(0.25),
                        y_cursor,
                        box_width - PInches(0.5),
                        remaining_h(),
                    )
                    if pic is not None:
                        y_cursor += pic.height + PInches(0.15)
                y_cursor += PInches(0.10)
                continue

            tokens = split_text_and_math(cleaned_b)

            def _inline_math_to_text(latex_body: str) -> str:
                s = (latex_body or "").strip()
                if not s:
                    return ""

                # ✅ FIX 1: re.sub must include replacement string
                s = re.sub(r"\\left|\\right", "", s)

                # ✅ FIX 2: replace braces correctly
                s = s.replace("{", "").replace("}", "")

                # whitespace-ish latex spacing
                s = s.replace(r"\,", " ").replace(r"\;", " ").replace(r"\:", " ").replace(r"\ ", " ")

                # operators
                s = s.replace(r"\times", "×")
                s = s.replace(r"\cdot", "·")
                s = s.replace(r"\leq", "≤").replace(r"\geq", "≥")
                s = s.replace(r"\neq", "≠")
                s = s.replace(r"\pm", "±")

                # simple fractions
                s = re.sub(r"\\frac\s*\(\s*([^()]*)\s*\)\s*\(\s*([^()]*)\s*\)", r"\1/\2", s)
                s = re.sub(r"\\frac\s*\{\s*([^{}]+)\s*\}\s*\{\s*([^{}]+)\s*\}", r"\1/\2", s)

                # coordinate pairs like (-4,-7) must survive
                if re.fullmatch(r"\(\s*-?\d+(\.\d+)?\s*,\s*-?\d+(\.\d+)?\s*\)", s):
                    return re.sub(r"\s+", " ", s).strip()

                # If it’s “simple text math”, keep it as normal PPT text
                if re.fullmatch(r"[0-9A-Za-z\s\+\-\=\(\)\[\],\./\^<>≤≥×·±]*", s):
                    return s.strip()

                # too complex -> render as image
                return ""

            parts: List[str] = []
            for tok in tokens:
                kind = tok[0]
                seg = tok[1]
                mtype = tok[2] if (kind == "math" and len(tok) > 2) else None
                if kind == "text":
                    parts.append(seg)
                elif mtype == "inline":
                    t = _inline_math_to_text(seg)
                    if t:
                        parts.append(t)

            bullet_text = "".join(parts).strip() or "(equation)"
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet_text
            p.level = 0
            p.font.size = Pt(20)

            eqs = [tok for tok in tokens if tok[0] == "math"]
            for _, latex_body, m_kind in eqs:
                # If inline math was safely converted to text, don't add an equation image
                if m_kind == "inline" and _inline_math_to_text(latex_body):
                    continue

                buf = render_latex_png(latex_body)
                img = Image.open(buf)
                tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                img.save(tmp_img.name)
                tmp_img.close()

                pic = _add_picture_fit(
                    slide,
                    tmp_img.name,
                    box_left + PInches(0.25),
                    y_cursor,
                    box_width - PInches(0.5),
                    min(remaining_h(), PInches(1.3 if m_kind == "display" else 1.0)),
                )
                if pic is not None:
                    y_cursor += pic.height + PInches(0.12)

            if visual_imgs:
                for img_path in visual_imgs:
                    pic = _add_picture_fit(
                        slide,
                        img_path,
                        box_left + PInches(0.25),
                        y_cursor,
                        box_width - PInches(0.5),
                        remaining_h(),
                    )
                    if pic is not None:
                        y_cursor += pic.height + PInches(0.15)

            y_cursor += PInches(0.10)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(out_path)


def build_mathjax_html(markdown_text: str) -> str:
    markdown_text = _strip_bundle_sections_for_docx(markdown_text)

    escaped = (markdown_text
               .replace("&", "&amp;")
               .replace("<", "&lt;")
               .replace(">", "&gt;"))
    escaped = escaped.replace("\n", "<br>")

    return f"""
<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8"/>
  <script>
    window.MathJax = {{
      tex: {{
        inlineMath: [['\\\\(','\\\\)']],
        displayMath: [['\\\\[','\\\\]']]
      }}
    }};
  </script>
  <script defer src="https://cdn.jsdelivr.net/npm/mathjax@3/es5/tex-chtml.js"></script>
  <style>
    body {{ font-family: Arial, sans-serif; padding: 20px; line-height: 1.35; }}
    .box {{ max-width: 980px; margin: 0 auto; }}
    .hint {{ color: #888; font-size: 12px; margin-bottom: 12px; }}
  </style>
</head>
<body>
  <div class="box">
    <div class="hint">Preview only. Edit the Markdown / PPT outline in the app, then export DOCX/PPTX.</div>
    {escaped}
  </div>
</body>
</html>
""".strip()
