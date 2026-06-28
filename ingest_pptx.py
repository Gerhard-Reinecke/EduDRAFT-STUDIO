# ======================================================================================
# ingest_pptx.py
# ======================================================================================
# Module: PPTX Ingestion & Draft Reconstruction Engine
#
# System: EduDraft Studio (Marike App)
# Version: 4.4.10
#
# --------------------------------------------------------------------------------------
# PURPOSE
# --------------------------------------------------------------------------------------
# This module handles the ingestion, extraction, and conversion of teacher-uploaded
# PPTX files into structured, editable educational drafts.
#
# It extracts readable slide content from presentation files, reconstructs that
# content into Markdown with LaTeX through the LLM pipeline, and saves the result
# into the Drafts system with automatic versioning and preview support.
#
# --------------------------------------------------------------------------------------
# CORE RESPONSIBILITIES
# --------------------------------------------------------------------------------------
# 1. PPTX Text Extraction
#    - Reads uploaded .pptx files using python-pptx
#    - Extracts slide text frames, bullet content, and table text
#    - Preserves slide sequencing for downstream reconstruction
#    - Applies character limits to protect token usage
#
# 2. Upload Analysis
#    - Validates uploaded PPTX files and file paths
#    - Reports file size and extraction results to the UI
#    - Enables or disables generation depending on readability
#
# 3. Draft Reconstruction (LLM Integration)
#    - Converts slide deck text into an editable classroom-ready draft
#    - Preserves slide flow, titles, bullet hierarchy, and lesson sequence where possible
#    - Supports optional memo / answer key generation
#
# 4. Draft Persistence
#    - Creates a new Draft record in Supabase
#    - Auto-saves Version 1 into draft_versions
#    - Stores contextual metadata such as subject, year level, and curriculum stream
#
# 5. Preview Generation
#    - Builds an HTML preview using MathJax-compatible rendering
#    - Supports immediate review of the generated draft output
#
# --------------------------------------------------------------------------------------
# DEPENDENCIES
# --------------------------------------------------------------------------------------
# - python-pptx → PPTX slide parsing
# - llm.py → Prompt construction and LLM interaction
# - exports.py → HTML preview rendering
# - auth.py → Session validation
# - rate_limit.py → Usage control
# - config.py → Supabase client
# - Gradio → UI interaction layer
#
# --------------------------------------------------------------------------------------
# DESIGN PRINCIPLES
# --------------------------------------------------------------------------------------
# - Teacher-first workflow with minimal friction
# - Best-effort content recovery from structured slide decks
# - Safe token handling through controlled truncation
# - Clean separation between ingestion, generation, preview, and persistence
# - Functional symmetry with PDF and DOCX ingestion workflows
#
# --------------------------------------------------------------------------------------
# NOTES
# --------------------------------------------------------------------------------------
# - Only .pptx files are supported
# - Extraction is best-effort and depends on text being present in slide objects
# - Image-heavy decks, embedded graphics, or non-standard content may yield limited text
# - OCR is not performed in this module
# - Large extracted outputs are truncated to preserve LLM performance
# - This module mirrors PDF/DOCX ingestion behavior for workflow consistency
#
# ======================================================================================

import os
import re
import uuid
import tempfile
from datetime import datetime, timezone

import gradio as gr

from llm import build_user_request, call_llm, split_sections, combine_doc_and_memo
from exports import build_mathjax_html
from auth import _require_session
from rate_limit import check_rate_limit
from config import supabase


# =============================
# PPTX INGESTION
# =============================
def extract_text_from_pptx(pptx_path: str, max_chars: int = 120_000) -> str:
    """
    Extracts text from a PPTX (slide text boxes + tables). Best-effort.
    Truncates to max_chars to protect token usage.
    """
    if not pptx_path or not os.path.exists(pptx_path):
        return ""

    text_runs = []

    try:
        from pptx import Presentation  # python-pptx
        prs = Presentation(pptx_path)

        for i, slide in enumerate(prs.slides, start=1):
            slide_bits = [f"Slide {i}:"]
            for shape in slide.shapes:
                # Text frames
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = (p.text or "").strip()
                        if t:
                            slide_bits.append(t)

                # Tables
                if hasattr(shape, "has_table") and shape.has_table:
                    tbl = shape.table
                    for row in tbl.rows:
                        row_cells = []
                        for cell in row.cells:
                            ct = re.sub(r"\s+", " ", (cell.text or "").strip())
                            if ct:
                                row_cells.append(ct)
                        if row_cells:
                            slide_bits.append(" | ".join(row_cells))

            # only append if we got meaningful content
            cleaned = "\n".join([b for b in slide_bits if b.strip()])
            if cleaned.strip() != f"Slide {i}:":
                text_runs.append(cleaned)

    except Exception:
        return ""

    text = "\n\n".join(text_runs).strip()
    text = re.sub(r"\n{3,}", "\n\n", text)

    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n[TRUNCATED]"

    return text


def action_analyze_pptx(pptx_file):
    """
    Called when a PPTX is uploaded.
    Returns:
      1) markdown status string
      2) gr.update(interactive=True/False)
    """
    try:
        if not pptx_file:
            return "No PPTX uploaded yet.", gr.update(interactive=False)

        if isinstance(pptx_file, dict):
            pptx_path = pptx_file.get("name") or pptx_file.get("path") or ""
            pptx_name = os.path.basename(pptx_path) if pptx_path else "uploaded.pptx"
        else:
            pptx_path = str(pptx_file)
            pptx_name = os.path.basename(pptx_path)

        if not pptx_path or not os.path.exists(pptx_path):
            return "⚠️ PPTX uploaded, but file path not found.", gr.update(interactive=False)

        size_kb = os.path.getsize(pptx_path) / 1024.0
        txt = extract_text_from_pptx(pptx_path)
        n = len(txt or "")

        if n == 0:
            msg = (
                f"📊 **PPTX detected**\n\n"
                f"- File: `{pptx_name}`\n"
                f"- Size: **{size_kb:.1f} KB**\n"
                f"- Extracted chars: **0** ❌\n\n"
                "This deck may be image-heavy or text could not be read."
            )
            return msg, gr.update(interactive=False)

        msg = (
            f"📊 **PPTX detected**\n\n"
            f"- File: `{pptx_name}`\n"
            f"- Size: **{size_kb:.1f} KB**\n"
            f"- Extracted chars: **{n}** ✅"
        )
        return msg, gr.update(interactive=True)

    except Exception as e:
        return f"⚠️ PPTX analysis failed: {type(e).__name__}: {e}", gr.update(interactive=False)


def pptx_to_draft_prompt(pptx_text: str, output_type: str, include_memo: bool) -> str:
    memo_line = "Yes" if include_memo else "No"
    return f"""
You are given text extracted from a teacher's PPTX. The teacher wants an EDITABLE draft.
Your job:
- Reconstruct the slide deck content into a clean, editable draft in Markdown with LaTeX.
- Preserve slide titles, bullet structure, and sequence as best as possible.
- If the deck is a lesson, keep lesson flow and pacing.
- If some content is unclear, make minimal reasonable assumptions and list 1–3 short assumptions at the end.
Target output type: {output_type}
Include answer key/memo: {memo_line}
PPTX_TEXT_START
<<<
{pptx_text}
>>>
PPTX_TEXT_END
Now produce the required output following the OUTPUT RULE exactly.
""".strip()


def action_generate_from_pptx(
    supabase_session,
    pptx_file,
    pptx_draft_title: str,
    education_level,
    country, state_province,
    uni_country, university_name, faculty, module_code,
    year_level, course, typed_subject, course_stream, output_type,
    include_memo, model_name
):
    """
    Upload PPTX -> generate NEW editable draft + auto-save v1.
    Mirrors PDF/DOCX behavior.
    """
    try:
        access_token, refresh_token, user_id, err = _require_session(supabase_session)
        if err:
            return "", "", "", None, f"❌ {err}", "", "PPTX upload", "", "", 0

        can_proceed, limit_msg = check_rate_limit(supabase_session, action="generate")
        if not can_proceed:
            return "", "", "", None, limit_msg, "", "PPTX upload", "", "", 0

        if not pptx_file:
            return "", "", "", None, f"❌ Please upload a PPTX first. {limit_msg}", "", "PPTX upload", "", "", 0

        if isinstance(pptx_file, dict):
            pptx_path = pptx_file.get("name") or pptx_file.get("path") or ""
            pptx_name = os.path.basename(pptx_path) if pptx_path else "uploaded.pptx"
        else:
            pptx_path = str(pptx_file)
            pptx_name = os.path.basename(pptx_path)

        if not pptx_path or not os.path.exists(pptx_path):
            return "", "", "", None, f"❌ Uploaded PPTX path not found. {limit_msg}", "", "PPTX upload", "", "", 0

        if os.path.splitext(pptx_name)[1].lower() != ".pptx":
            return "", "", "", None, f"❌ That file is not a PPTX. {limit_msg}", "", "PPTX upload", "", "", 0

        pptx_text = extract_text_from_pptx(pptx_path)
        if not (pptx_text or "").strip():
            return "", "", "", None, (
                "❌ I couldn't extract any readable text from this PPTX.\n\n"
                f"File: {pptx_name}\n"
                f"Extracted chars: {len(pptx_text or '')}\n\n"
                f"{limit_msg}"
            ), "", "PPTX upload", "", "", 0

        # -----------------------------
        # SUBJECT (GLOBAL, NOT AU-ONLY)
        # -----------------------------
        if education_level == "University / Tertiary":
            eff_country = (uni_country or "").strip() or "Not specified"
            eff_state = ""
            effective_course = (module_code or "").strip() or "Not specified"
        else:
            eff_country = (country or "").strip() or "Not specified"
            eff_state = (state_province or "").strip()

            chosen = (course or "").strip()
            if chosen == "Other (type it)":
                ts = (typed_subject or "").strip()
                if not ts:
                    return "", "", "", None, f"❌ Subject is required. {limit_msg}", "", "PPTX upload", "", "", 0
                effective_course = ts
            else:
                effective_course = chosen or "Not specified"

        curriculum_stream = (course_stream or "").strip()

        instruction_text = f"Generate an editable draft from the uploaded PPTX: {pptx_name}"

        teacher_context = build_user_request(
            instruction_text=instruction_text,
            education_level=education_level,
            country=eff_country,
            state_province=eff_state,
            year_level=year_level,
            course=effective_course,
            output_type=output_type,
            include_memo=include_memo,
            university_name=university_name,
            faculty=faculty,
            module_code=module_code,
            curriculum_stream=curriculum_stream,
        )

        prompt = teacher_context + "\n\n" + pptx_to_draft_prompt(pptx_text, output_type, include_memo)

        llm_text = call_llm(prompt, model_name, edit_mode=False)
        doc_md, ppt_outline, answer_key = split_sections(llm_text)
        combined_md = combine_doc_and_memo(doc_md, answer_key, include_memo)

        tmpdir = tempfile.mkdtemp()
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        html_path = os.path.join(tmpdir, f"Preview_{stamp}.html")
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(build_mathjax_html(combined_md))

        now = datetime.now(timezone.utc).isoformat()
        title = (pptx_draft_title or "").strip()
        if not title:
            title = os.path.splitext(pptx_name)[0].strip() or "PPTX Draft"

        draft_id = str(uuid.uuid4())

        draft_row = {
            "id": draft_id,
            "user_id": user_id,
            "title": title,
            "subject": (effective_course or "").strip() or None,
            "curriculum_stream": (curriculum_stream or "").strip() or None,
            "education_level": (education_level or "").strip() or None,
            "country": (eff_country or "").strip() or None,
            "state_province": (eff_state or "").strip() or None,
            "year_level": (year_level or "").strip() or None,
            "course": (effective_course or "").strip() or None,
            "created_at": now,
            "updated_at": now,
        }
        supabase.table("drafts").insert(draft_row).execute()

        version_row = {
            "draft_id": draft_id,
            "version": 1,
            "doc_md": combined_md or "",
            "ppt_outline": ppt_outline or "",
            "created_at": now,
        }
        supabase.table("draft_versions").insert(version_row).execute()

        status = (
            "✅ PPTX converted into a NEW editable draft and auto-saved as Version 1.\n"
            f"Draft name: {title}\n"
            f"{limit_msg}\n\n"
            "You can now edit it and Save NEW versions anytime."
        )

        return (
            f"PPTX used: {pptx_name}",
            combined_md,
            ppt_outline,
            html_path,
            status,
            instruction_text,
            "PPTX upload",
            f"PPTX used: {pptx_name}",
            draft_id,
            1
        )

    except Exception as e:
        from llm import safe_err
        return "", "", "", None, safe_err("PPTX → Draft failed.", e), "", "PPTX upload", "", "", 0